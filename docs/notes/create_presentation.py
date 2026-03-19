"""Generate PowerPoint presentation: Unverified Connectivity Root Cause Analysis.

12-slide deck:
  1. Title
  2. The Problem (lb/ub doesn't encode correct point order)
  3. 2D: How lb/ub works (text)
  4. 2D Diagonals diagram
  5. From 2D to 3D diagram
  6. The Mismatch (combined traversal + mismatch diagram)
  7. The Fix (natural vs swapped diagram)
  8. The 8 Permutation Matrices (table)
  9. The Solution (orientation flags + permutation_index)
  10. Python Reference (NASA Plot3D_utilities + plot3d-rs)
  11. Summary
  12. JSON Output Example

Run: /opt/homebrew/Caskroom/miniconda/base/envs/dev/bin/python create_presentation.py
"""
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_MED = RGBColor(0x24, 0x24, 0x3E)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED = RGBColor(0xFF, 0x66, 0x66)
GREEN = RGBColor(0x66, 0xFF, 0x99)
YELLOW = RGBColor(0xFF, 0xDD, 0x57)
BLUE = RGBColor(0x66, 0xBB, 0xFF)
CODE_BG = RGBColor(0x18, 0x18, 0x28)

MPL_BG = '#1B1B2F'
MPL_FACE1 = '#66BBFF'
MPL_FACE2 = '#FF6666'
MPL_MATCH = '#66FF99'
MPL_ACCENT = '#4EC9B0'
MPL_YELLOW = '#FFDD57'
MPL_GRAY = '#CCCCCC'


def fig_to_image_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight',
                facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = LIGHT_GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)


def _add_title_bar(slide, title):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_MED
    rect.line.fill.background()
    tf = rect.text_frame
    tf.margin_left = Inches(0.8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT


def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    _add_title_bar(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level, color, size = item
        else:
            text, level, color, size = item, 0, WHITE, 22
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = level
        p.space_before = Pt(6) if level > 0 else Pt(14)


def add_table_slide(title, headers, rows, col_widths=None, footer_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    _add_title_bar(slide, title)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5) * n_rows)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = LIGHT_GRAY
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_MED if i % 2 == 0 else RGBColor(0x2A, 0x2A, 0x48)
    if footer_text:
        y_bottom = 1.6 + 0.5 * n_rows + 0.3
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y_bottom), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = footer_text
        p.font.size = Pt(20)
        p.font.color.rgb = YELLOW


def add_image_slide(title, image_stream, subtitle=""):
    from PIL import Image as PILImage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    _add_title_bar(slide, title)
    image_stream.seek(0)
    pil_img = PILImage.open(image_stream)
    img_w, img_h = pil_img.size
    image_stream.seek(0)
    img_aspect = img_w / img_h
    max_w, max_h = 12.3, 5.6
    area_aspect = max_w / max_h
    if img_aspect > area_aspect:
        w = max_w
        h = w / img_aspect
    else:
        h = max_h
        w = h * img_aspect
    left = (13.333 - w) / 2
    top = 1.2 + (max_h - h) / 2
    slide.shapes.add_picture(image_stream, Inches(left), Inches(top), Inches(w), Inches(h))
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = YELLOW
        p.alignment = PP_ALIGN.CENTER


def _setup_ax(ax):
    ax.set_facecolor(MPL_BG)
    ax.tick_params(colors=MPL_GRAY)
    for spine in ax.spines.values():
        spine.set_color(MPL_GRAY)


# ══════════════════════════════════════════════════════════════════════
# DIAGRAM GENERATORS (only the 4 we need)
# ══════════════════════════════════════════════════════════════════════

def make_2d_diagonals_diagram():
    """4 diagonal lb/ub encodings on a 2D grid. All green = all encodable."""
    fig, ax = plt.subplots(figsize=(12, 6), facecolor=MPL_BG)
    _setup_ax(ax)
    ni, nj = 4, 5
    for i in range(ni):
        for j in range(nj):
            ax.plot(j, i, 'o', color='#555588', markersize=12, zorder=2)
    diags = [
        ((0, 0), (ni-1, nj-1), 'lb=[0,0]  ub=[3,4]\ni fwd, j fwd', '#66FF99', (0.25, 0.15)),
        ((ni-1, nj-1), (0, 0), 'lb=[3,4]  ub=[0,0]\ni rev, j rev', '#4EC9B0', (-0.25, -0.15)),
        ((0, nj-1), (ni-1, 0), 'lb=[0,4]  ub=[3,0]\ni fwd, j rev', '#66BBFF', (0.25, -0.15)),
        ((ni-1, 0), (0, nj-1), 'lb=[3,0]  ub=[0,4]\ni rev, j fwd', '#FFDD57', (-0.25, 0.15)),
    ]
    for (i0, j0), (i1, j1), label, color, (dy, dx) in diags:
        ax.annotate('', xy=(j1+dx*0.5, i1+dy*0.5), xytext=(j0+dx*0.5, i0+dy*0.5),
                    arrowprops=dict(arrowstyle='->', color=color, lw=3, alpha=0.85))
        ax.plot(j0+dx*0.5, i0+dy*0.5, 'o', color=color, markersize=10, zorder=5)
        ax.plot(j1+dx*0.5, i1+dy*0.5, 'o', color=color, markersize=10,
                markerfacecolor='none', markeredgewidth=2.5, zorder=5)
        mid_j = (j0+j1)/2 + dx*2.5
        mid_i = (i0+i1)/2 + dy*2.5
        ax.text(mid_j, mid_i, label, color=color, fontsize=9, ha='center', va='center',
                fontweight='bold', family='monospace',
                bbox=dict(boxstyle='round,pad=0.3', facecolor=MPL_BG, edgecolor=color, alpha=0.9))
    ax.set_xlim(-1.5, nj+1.5)
    ax.set_ylim(-1.8, ni+1.5)
    ax.set_aspect('equal')
    ax.set_xticks(range(nj))
    ax.set_xticklabels([f'j={j}' for j in range(nj)], fontsize=10, color=MPL_GRAY)
    ax.set_yticks(range(ni))
    ax.set_yticklabels([f'i={i}' for i in range(ni)], fontsize=10, color=MPL_GRAY)
    ax.set_xlabel("j axis", color=MPL_GRAY, fontsize=12)
    ax.set_ylabel("i axis", color=MPL_GRAY, fontsize=12)
    ax.set_title("2D Face: 4 Diagonal Directions — All Encodable by lb/ub",
                 color='white', fontsize=15, fontweight='bold', pad=12)
    ax.text(nj/2-0.5, -1.5,
            "2 axes  x  2 directions each  =  4 combinations\n"
            "Filled dot = lb (start)     Open dot = ub (end)",
            ha='center', fontsize=11, color=MPL_GRAY,
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#2A2A48', edgecolor=MPL_GRAY, alpha=0.8))
    plt.tight_layout()
    return fig_to_image_stream(fig)


def make_2d_to_3d_diagram():
    """Same-axis (4 combos) vs cross-axis (8 combos)."""
    fig, axes = plt.subplots(1, 2, figsize=(15, 6.5), facecolor=MPL_BG)
    for ax in axes:
        ax.set_facecolor(MPL_BG)
        ax.axis('off')

    # Left: same constant axis
    ax = axes[0]
    ax.set_xlim(0, 10); ax.set_ylim(-0.5, 9)
    ax.set_title("Same Constant Axis\n(both K-constant)", color=MPL_MATCH, fontsize=14,
                 fontweight='bold', pad=10)
    box_a = dict(boxstyle='round,pad=0.5', facecolor='#1A3A1A', edgecolor=MPL_MATCH, linewidth=2)
    ax.text(2.5, 7.5, "Face A (K=0)\naxes: i, j", fontsize=12, color=MPL_MATCH,
            ha='center', fontweight='bold', bbox=box_a)
    ax.text(7.5, 7.5, "Face B (K=0)\naxes: i, j", fontsize=12, color=MPL_MATCH,
            ha='center', fontweight='bold', bbox=box_a)
    labels = ["i fwd, j fwd", "i fwd, j rev", "i rev, j fwd", "i rev, j rev"]
    for idx, label in enumerate(labels):
        y = 5.5 - idx * 1.1
        ax.annotate('', xy=(6.2, y), xytext=(3.8, y),
                    arrowprops=dict(arrowstyle='->', color=MPL_MATCH, lw=2))
        ax.text(5.0, y+0.3, label, fontsize=9, color=MPL_MATCH, ha='center', family='monospace')
    ax.text(5.0, 0.8, "4 combinations\nAll encodable by lb/ub", fontsize=13, color=MPL_MATCH,
            ha='center', fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#114411', edgecolor=MPL_MATCH))

    # Right: different constant axes
    ax = axes[1]
    ax.set_xlim(0, 10); ax.set_ylim(-2.5, 9)
    ax.set_title("Different Constant Axes\n(K-constant meets J-constant)", color=MPL_YELLOW,
                 fontsize=14, fontweight='bold', pad=10)
    box_a2 = dict(boxstyle='round,pad=0.5', facecolor='#1A2A3A', edgecolor=MPL_FACE1, linewidth=2)
    ax.text(2.0, 7.5, "Face A (K=0)\naxes: i, j", fontsize=12, color=MPL_FACE1,
            ha='center', fontweight='bold', bbox=box_a2)
    box_b2 = dict(boxstyle='round,pad=0.5', facecolor='#3A1A1A', edgecolor=MPL_FACE2, linewidth=2)
    ax.text(8.0, 7.5, "Face B (J=0)\naxes: i, k", fontsize=12, color=MPL_FACE2,
            ha='center', fontweight='bold', bbox=box_b2)

    ax.text(5.0, 6.4, "No swap: A(i,j) maps to B(i,k)", fontsize=10, color=MPL_MATCH,
            ha='center', fontweight='bold')
    # -> = forward, <- = reversed
    no_swap = [
        ("i->i,  j->k", True, True),    # both fwd
        ("i->i,  j<-k", True, False),   # i fwd, k rev
        ("i<-i,  j->k", False, True),   # i rev, k fwd
        ("i<-i,  j<-k", False, False),  # both rev
    ]
    for idx, (label, i_fwd, k_fwd) in enumerate(no_swap):
        y = 5.6 - idx * 0.7
        ax.annotate('', xy=(6.8, y), xytext=(3.2, y),
                    arrowprops=dict(arrowstyle='->', color=MPL_MATCH, lw=1.5, alpha=0.7))
        ax.text(5.0, y+0.2, label, fontsize=8, color=MPL_MATCH, ha='center', family='monospace')

    ax.text(5.0, 2.9, "Swapped: A(i,j) maps to B(k,i)", fontsize=10, color=MPL_YELLOW,
            ha='center', fontweight='bold')
    swap = [
        ("i->k,  j->i", True, True),    # both fwd
        ("i->k,  j<-i", True, False),   # i->k fwd, j->i rev
        ("i<-k,  j->i", False, True),   # i->k rev, j->i fwd
        ("i<-k,  j<-i", False, False),  # both rev
    ]
    for idx, (label, _, _) in enumerate(swap):
        y = 2.2 - idx * 0.7
        ax.annotate('', xy=(6.8, y), xytext=(3.2, y),
                    arrowprops=dict(arrowstyle='->', color=MPL_YELLOW, lw=1.5, alpha=0.7))
        ax.text(5.0, y+0.2, label, fontsize=8, color=MPL_YELLOW, ha='center', family='monospace')

    ax.text(9.2, 4.5, "4", fontsize=16, color=MPL_MATCH, ha='center', fontweight='bold')
    ax.text(9.2, 4.0, "lb/ub", fontsize=9, color=MPL_MATCH, ha='center')
    ax.text(9.2, 1.2, "4", fontsize=16, color=MPL_YELLOW, ha='center', fontweight='bold')
    ax.text(9.2, 0.7, "need\nswap", fontsize=9, color=MPL_YELLOW, ha='center')
    ax.text(5.0, -2.0, "8 combinations total\nlb/ub only encodes 4 — need 'swapped' flag for the other 4",
            fontsize=11, color=MPL_FACE2, ha='center', fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#441111', edgecolor=MPL_FACE2))

    fig.suptitle("From 2D to 3D: Why 4 Combinations Become 8",
                 color='white', fontsize=16, fontweight='bold', y=1.0)
    plt.tight_layout()
    return fig_to_image_stream(fig)


def make_combined_mismatch_diagram():
    """Combined diagram: grid traversal (top) + mismatch explanation (bottom annotation).
    Shows K-const face (i outer, j inner) vs J-const face (i outer, k inner)
    with numbered points and highlights that strip sizes differ."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6), facecolor=MPL_BG)
    for ax in axes:
        _setup_ax(ax)

    # --- Left: K-constant face (outer=i, inner=j) ---
    ax = axes[0]
    ni, nj = 4, 6
    ax.set_title("Face A — K-constant\nfor i (outer), for j (inner)",
                 color=MPL_FACE1, fontsize=13, fontweight='bold', pad=10)
    ax.set_xlabel("j axis (inner — varies fastest)", color=MPL_GRAY, fontsize=11)
    ax.set_ylabel("i axis (outer)", color=MPL_GRAY, fontsize=11)

    point_num = 0
    coords = []
    for i in range(ni):
        for j in range(nj):
            coords.append((j, ni - 1 - i))
            # Highlight first 5 points to track them
            if point_num < 5:
                ax.plot(j, ni-1-i, 's', color=MPL_MATCH, markersize=28, zorder=2, alpha=0.3)
            ax.plot(j, ni-1-i, 'o', color=MPL_FACE1, markersize=18, zorder=3)
            ax.text(j, ni-1-i, str(point_num), color='white', fontsize=8,
                    ha='center', va='center', fontweight='bold', zorder=4)
            point_num += 1

    for idx in range(len(coords) - 1):
        x0, y0 = coords[idx]
        x1, y1 = coords[idx+1]
        dx, dy = x1-x0, y1-y0
        ax.annotate('', xy=(x1-0.15*dx, y1-0.15*dy), xytext=(x0+0.15*dx, y0+0.15*dy),
                    arrowprops=dict(arrowstyle='->', color=MPL_FACE1, lw=1.2, alpha=0.5))

    # Strip highlight
    rect = mpatches.FancyBboxPatch((-0.4, ni-1-0.4), nj-0.2, 0.8,
                                    boxstyle="round,pad=0.1", linewidth=2,
                                    edgecolor=MPL_YELLOW, facecolor='none', linestyle='--')
    ax.add_patch(rect)
    ax.text(nj/2-0.5, ni+0.3, "strip = 6 pts (j varies)", color=MPL_YELLOW, fontsize=9, ha='center')

    ax.set_xlim(-0.8, nj-0.2)
    ax.set_ylim(-0.8, ni+0.8)
    ax.set_aspect('equal')
    ax.set_xticks(range(nj))
    ax.set_xticklabels([f'j={j}' for j in range(nj)], fontsize=8, color=MPL_GRAY)
    ax.set_yticks(range(ni))
    ax.set_yticklabels([f'i={ni-1-i}' for i in range(ni)], fontsize=8, color=MPL_GRAY)

    # --- Right: J-constant face (outer=i, inner=k) ---
    ax = axes[1]
    ni2, nk = 6, 4
    ax.set_title("Face B — J-constant\nfor i (outer), for k (inner)",
                 color=MPL_FACE2, fontsize=13, fontweight='bold', pad=10)
    ax.set_xlabel("k axis (inner — varies fastest)", color=MPL_GRAY, fontsize=11)
    ax.set_ylabel("i axis (outer)", color=MPL_GRAY, fontsize=11)

    point_num = 0
    coords2 = []
    for i in range(ni2):
        for k in range(nk):
            coords2.append((k, ni2-1-i))
            if point_num < 5:
                ax.plot(k, ni2-1-i, 's', color=MPL_MATCH, markersize=28, zorder=2, alpha=0.3)
            ax.plot(k, ni2-1-i, 'o', color=MPL_FACE2, markersize=18, zorder=3)
            ax.text(k, ni2-1-i, str(point_num), color='white', fontsize=8,
                    ha='center', va='center', fontweight='bold', zorder=4)
            point_num += 1

    for idx in range(len(coords2) - 1):
        x0, y0 = coords2[idx]
        x1, y1 = coords2[idx+1]
        dx, dy = x1-x0, y1-y0
        ax.annotate('', xy=(x1-0.15*dx, y1-0.15*dy), xytext=(x0+0.15*dx, y0+0.15*dy),
                    arrowprops=dict(arrowstyle='->', color=MPL_FACE2, lw=1.2, alpha=0.5))

    rect2 = mpatches.FancyBboxPatch((-0.4, ni2-1-0.4), nk-0.2, 0.8,
                                     boxstyle="round,pad=0.1", linewidth=2,
                                     edgecolor=MPL_YELLOW, facecolor='none', linestyle='--')
    ax.add_patch(rect2)
    ax.text(nk/2-0.5, ni2+0.3, "strip = 4 pts (k varies)", color=MPL_YELLOW, fontsize=9, ha='center')

    ax.set_xlim(-0.8, nk-0.2)
    ax.set_ylim(-0.8, ni2+0.8)
    ax.set_aspect('equal')
    ax.set_xticks(range(nk))
    ax.set_xticklabels([f'k={k}' for k in range(nk)], fontsize=8, color=MPL_GRAY)
    ax.set_yticks(range(ni2))
    ax.set_yticklabels([f'i={ni2-1-i}' for i in range(ni2)], fontsize=8, color=MPL_GRAY)

    fig.text(0.5, -0.05,
             "Point 1 on Face A = (i=0, j=1)  but  Point 1 on Face B = (i=5, k=1)\n"
             "Strip sizes differ (6 vs 4) so the point sequences go out of sync!",
             ha='center', fontsize=12, color=MPL_YELLOW,
             bbox=dict(boxstyle='round,pad=0.5', facecolor='#2A2A48', edgecolor=MPL_YELLOW, alpha=0.9))

    fig.subplots_adjust(top=0.92, bottom=0.15)
    return fig_to_image_stream(fig)


def make_natural_vs_swapped_diagram():
    """Natural order (FAIL) vs swapped order (PASS) for Face B."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6.5), facecolor=MPL_BG)
    for ax in axes:
        _setup_ax(ax)

    ni, nk = 5, 4

    # Left: Natural order (FAILS)
    ax = axes[0]
    ax.set_title("Face B — NATURAL order\nfor i (outer), for k (inner)",
                 color=MPL_FACE2, fontsize=13, fontweight='bold', pad=8)
    point_num = 0
    coords = []
    for i in range(ni):
        for k in range(nk):
            coords.append((k, ni-1-i))
            ax.plot(k, ni-1-i, 'o', color=MPL_FACE2, markersize=24, zorder=3)
            ax.text(k, ni-1-i, str(point_num), color='white', fontsize=9,
                    ha='center', va='center', fontweight='bold', zorder=4)
            point_num += 1
    for idx in range(len(coords) - 1):
        x0, y0 = coords[idx]
        x1, y1 = coords[idx+1]
        dx, dy = x1-x0, y1-y0
        ax.annotate('', xy=(x1-0.18*dx, y1-0.18*dy), xytext=(x0+0.18*dx, y0+0.18*dy),
                    arrowprops=dict(arrowstyle='->', color=MPL_FACE2, lw=1.5, alpha=0.6))
    ax.set_xlabel("k (inner)", color=MPL_GRAY, fontsize=11)
    ax.set_ylabel("i (outer)", color=MPL_GRAY, fontsize=11)
    ax.set_xlim(-0.7, nk-0.3); ax.set_ylim(-1.5, ni-0.3)
    ax.set_aspect('equal')
    ax.set_xticks(range(nk)); ax.set_yticks(range(ni))
    ax.set_yticklabels([f'i={ni-1-y}' for y in range(ni)], fontsize=9, color=MPL_GRAY)
    ax.set_xticklabels([f'k={x}' for x in range(nk)], fontsize=9, color=MPL_GRAY)
    ax.text(nk/2-0.5, -1.2, "FAIL: strip = 4 (doesn't match Face A's strip of 5)",
            color=MPL_FACE2, fontsize=11, ha='center', fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#441111', edgecolor=MPL_FACE2))

    # Right: Swapped order (PASSES)
    ax = axes[1]
    ax.set_title("Face B — SWAPPED order\nfor k (outer), for i reversed (inner)",
                 color=MPL_MATCH, fontsize=13, fontweight='bold', pad=8)
    point_num = 0
    coords2 = []
    for k in range(nk):
        for i in range(ni-1, -1, -1):
            coords2.append((k, i))
            ax.plot(k, i, 'o', color=MPL_MATCH, markersize=24, zorder=3)
            ax.text(k, i, str(point_num), color='#1B1B2F', fontsize=9,
                    ha='center', va='center', fontweight='bold', zorder=4)
            point_num += 1
    for idx in range(len(coords2) - 1):
        x0, y0 = coords2[idx]
        x1, y1 = coords2[idx+1]
        dx, dy = x1-x0, y1-y0
        if abs(dx) > 0.01 or abs(dy) > 0.01:
            ax.annotate('', xy=(x1-0.18*dx, y1-0.18*dy), xytext=(x0+0.18*dx, y0+0.18*dy),
                        arrowprops=dict(arrowstyle='->', color=MPL_MATCH, lw=1.5, alpha=0.6))
    ax.set_xlabel("k (now OUTER — swapped!)", color=MPL_YELLOW, fontsize=11, fontweight='bold')
    ax.set_ylabel("i reversed (now INNER)", color=MPL_YELLOW, fontsize=11, fontweight='bold')
    ax.set_xlim(-0.7, nk-0.3); ax.set_ylim(-1.5, ni-0.3)
    ax.set_aspect('equal')
    ax.set_xticks(range(nk)); ax.set_yticks(range(ni))
    ax.set_yticklabels([f'i={y}' for y in range(ni)], fontsize=9, color=MPL_GRAY)
    ax.set_xticklabels([f'k={x}' for x in range(nk)], fontsize=9, color=MPL_GRAY)
    ax.text(nk/2-0.5, -1.2, "PASS: strip = 5 — matches Face A!",
            color=MPL_MATCH, fontsize=11, ha='center', fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#114411', edgecolor=MPL_MATCH))

    fig.subplots_adjust(top=0.92)
    return fig_to_image_stream(fig)


# ══════════════════════════════════════════════════════════════════════
# BUILD SLIDES (9 total)
# ══════════════════════════════════════════════════════════════════════
print("Generating presentation (12 slides)...")

# ── Slide 1: Title ──
add_title_slide(
    "Unverified Connectivity Faces",
    "Root Cause Analysis & Solution"
)

# ── Slide 2: The Problem ──
add_table_slide(
    "The Problem: lb/ub Doesn't Encode the Correct Point Order",
    ["Category", "Count", "Status"],
    [
        ["connectivity_face_matches", "14,389", "Verified point-by-point"],
        ["connectivity_unverified", "10,469", "Geometrically matched, failed directed verification"],
        ["non_connected_faces", "1,886", "No match found (boundary faces)"],
    ],
    col_widths=[4.0, 1.5, 6.2],
    footer_text="lb/ub works for same-plane connections but fails when faces are on different constant planes.",
)

# ── Slide 3: 2D — How lb/ub Works ──
add_content_slide(
    "Start Simple: 2D Grid Face Connectivity",
    [
        "A face on a structured grid is a 2D surface defined by two varying axes",
        ("Example: a K-constant face varies in i and j", 1, LIGHT_GRAY, 20),
        "",
        "The lb/ub convention encodes the diagonal — start corner and end corner",
        ("lb = lower bound (start),   ub = upper bound (end)", 1, BLUE, 20),
        "",
        "For two axes (i, j), each can go forward or backward:",
        ("i: forward (0 -> 24) or backward (24 -> 0)", 1, LIGHT_GRAY, 18),
        ("j: forward (0 -> 408) or backward (408 -> 0)", 1, LIGHT_GRAY, 18),
        "",
        ("2 directions x 2 axes = 4 diagonal combinations — all encodable by lb/ub", 0, GREEN, 22),
    ],
)

# ── Slide 4: 2D Diagonals Diagram ──
print("  [1/4] 2D diagonals diagram...")
img_2d = make_2d_diagonals_diagram()
add_image_slide(
    "2D Diagonals: 4 Ways to Traverse a Face",
    img_2d,
    "Each arrow = one lb/ub encoding. All 4 directions are representable."
)

# ── Slide 5: From 2D to 3D ──
print("  [2/4] 2D to 3D diagram...")
img_2d3d = make_2d_to_3d_diagram()
add_image_slide(
    "From 2D to 3D: When 4 Combinations Become 8",
    img_2d3d,
    "Same constant axis = 4 combos (lb/ub works).  Different axes = 8 combos (need swap flag)."
)

# ── Slide 6: The Mismatch (combined) ──
print("  [3/4] Combined mismatch diagram...")
img_mismatch = make_combined_mismatch_diagram()
add_image_slide(
    "Why It Fails: Different Strip Sizes = Wrong Point Order",
    img_mismatch,
    "Face A numbers points in strips of 6 (j inner), Face B in strips of 4 (k inner) — they don't align."
)

# ── Slide 7: The Fix ──
print("  [4/4] Natural vs swapped diagram...")
img_fix = make_natural_vs_swapped_diagram()
add_image_slide(
    "The Fix: Swap the Loop Order So Points Match",
    img_fix,
    "swapped=true, v_reversed=true makes Face B's point sequence match Face A exactly."
)

# ── Slide 8: The 8 Permutation Matrices ──
add_table_slide(
    "The 8 Permutation Matrices",
    ["Index", "Bits", "u_rev", "v_rev", "swapped", "Matrix", "Description"],
    [
        ["0", "000", "No",  "No",  "No",  "[[1,0],[0,1]]",   "Identity"],
        ["1", "001", "Yes", "No",  "No",  "[[-1,0],[0,1]]",  "u reversed"],
        ["2", "010", "No",  "Yes", "No",  "[[1,0],[0,-1]]",  "v reversed"],
        ["3", "011", "Yes", "Yes", "No",  "[[-1,0],[0,-1]]", "Both reversed"],
        ["4", "100", "No",  "No",  "Yes", "[[0,1],[1,0]]",   "Swapped"],
        ["5", "101", "Yes", "No",  "Yes", "[[0,-1],[1,0]]",  "Swap + u rev"],
        ["6", "110", "No",  "Yes", "Yes", "[[0,1],[-1,0]]",  "Swap + v rev"],
        ["7", "111", "Yes", "Yes", "Yes", "[[0,-1],[-1,0]]", "Swap + both"],
    ],
    col_widths=[0.8, 0.8, 1.0, 1.0, 1.2, 2.8, 3.0],
    footer_text="index = u_reversed | (v_reversed << 1) | (swapped << 2)",
)

# ── Slide 9: The Solution — Orientation + permutation_index ──
add_content_slide(
    "The Solution: Orientation with permutation_index",
    [
        "Both Python and Rust encode orientation as a permutation_index (0-7):",
        ("index = u_reversed | (v_reversed << 1) | (swapped << 2)", 1, ACCENT, 18),
        "",
        "Python connectivity() returns an orientation dict per face match:",
        ("{'permutation_index': 5, 'plane': 'cross-plane', 'permutation_matrix': [[-1,0],[1,0]]}", 1, GREEN, 14),
        "",
        "Rust plot3d-rs stores it as Orientation struct:",
        ("Orientation { permutation_index: u8, plane: OrientationPlane }", 1, GREEN, 15),
        "",
        "PERMUTATION_MATRICES constant exported from both libraries:",
        ("Python: from plot3d import PERMUTATION_MATRICES  # shape (8, 2, 2)", 1, YELLOW, 14),
        ("Rust: use plot3d::PERMUTATION_MATRICES;  // [[[i8; 2]; 2]; 8]", 1, YELLOW, 14),
        "",
        ("With permutation_index, all cross-plane faces pass verification", 0, GREEN, 22),
    ],
)

# ── Slide 10: Python Reference — Plot3D_utilities ──
add_content_slide(
    "Python & Rust Reference",
    [
        "Python (NASA Plot3D_utilities):",
        ("github.com/nasa/Plot3D_utilities", 1, BLUE, 18),
        ("connectivity_fast() returns orientation dict with permutation_index per match", 1, LIGHT_GRAY, 15),
        ("PERMUTATION_MATRICES exported — 8 canonical 2x2 signed permutation matrices", 1, LIGHT_GRAY, 15),
        "",
        "Rust (plot3d-rs):",
        ("github.com/pjuangph/plot3d-rs", 1, BLUE, 18),
        ("FaceMatch.orientation: Option<Orientation> with permutation_index: u8", 1, LIGHT_GRAY, 15),
        ("PERMUTATION_MATRICES constant: [[[i8; 2]; 2]; 8]", 1, LIGHT_GRAY, 15),
        "",
        "Both use the same bit encoding and 8-permutation system",
        ("Verification: try specified permutation first, fall back to brute-force all 8", 1, ACCENT, 15),
    ],
)

# ── Slide 11: Summary ──
add_content_slide(
    "Summary",
    [
        "10,469 unverified faces are geometrically correct connections",
        ("Blocks whose faces lie on different constant planes (K <-> J, etc.)", 1, YELLOW, 18),
        "",
        "Root cause: lb/ub encodes direction but not loop nesting order",
        ("2D: 4 diagonals — all encodable by lb/ub", 1, GREEN, 16),
        ("3D cross-plane: 8 permutations — lb/ub only encodes 4", 1, RED, 16),
        ("The 'swapped' flag covers the other 4", 1, YELLOW, 16),
        "",
        "The solution: permutation_index (0-7) encodes all 8 orientations",
        ("index = u_reversed | (v_reversed << 1) | (swapped << 2)", 1, ACCENT, 16),
        ("Both Python and Rust export PERMUTATION_MATRICES constant", 1, GREEN, 16),
    ],
)

# ── Slide 12: JSON Output Example ──
add_content_slide(
    "JSON Output: What Solvers Receive",
    [
        "Each face_match in the JSON output now includes orientation:",
        "",
        ("{", 1, ACCENT, 16),
        ('  "block1": {"block_index": 0, "lb": [0,0,0], "ub": [64,0,64]},', 1, LIGHT_GRAY, 14),
        ('  "block2": {"block_index": 421, "lb": [0,0,0], "ub": [0,64,64]},', 1, LIGHT_GRAY, 14),
        ('  "orientation": {', 1, YELLOW, 14),
        ('    "permutation_index": 5,', 1, YELLOW, 14),
        ('    "plane": "cross-plane"', 1, YELLOW, 14),
        ("  }", 1, YELLOW, 14),
        ("}", 1, ACCENT, 16),
        "",
        ("Solvers use permutation_index to look up PERMUTATION_MATRICES[5]", 0, GREEN, 18),
        ("and apply the correct 2x2 transform when exchanging boundary data", 0, GREEN, 18),
    ],
)

# ── Save ──
import os
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(SCRIPT_DIR, "cross_plane_connectivity_findings.pptx")
prs.save(out_path)
print(f"\nSaved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
